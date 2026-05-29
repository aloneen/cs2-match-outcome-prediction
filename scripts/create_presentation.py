"""create_presentation.py — CS2 defense slides, editorial / NotebookLM style"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Color palette ─────────────────────────────────────────────────
BLACK  = RGBColor(0x09, 0x09, 0x09)   # near-black slide bg
DGREY  = RGBColor(0x1E, 0x1E, 0x1E)   # dark card bg
MGREY  = RGBColor(0x2E, 0x2E, 0x2E)   # rule on dark
OFFWH  = RGBColor(0xF6, 0xF6, 0xF4)   # warm off-white slide bg
CARD   = RGBColor(0xEC, 0xEC, 0xEA)   # card on light slides
RULEW  = RGBColor(0xD4, 0xD4, 0xD2)   # rule on light slides
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LTXT   = RGBColor(0xE8, 0xE8, 0xE6)   # primary text on dark
LGREY  = RGBColor(0x88, 0x88, 0x86)   # secondary text on dark
DNUM   = RGBColor(0x20, 0x20, 0x20)   # ghost numbers on dark
DTXT   = RGBColor(0x18, 0x18, 0x18)   # primary text on light
STXT   = RGBColor(0x78, 0x78, 0x76)   # secondary text on light

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = 13.33, 7.5
BLANK = prs.slide_layouts[6]
FIG   = "report_figures"
FONT  = "Calibri"


# ── Primitives ────────────────────────────────────────────────────

def R(sl, x, y, w, h, c):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background()
    return s

def T(sl, x, y, w, h, text, sz=16, bold=False, italic=False,
      c=None, align=PP_ALIGN.LEFT):
    if c is None: c = DTXT
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = c
    return tb

def I(sl, path, x, y, w, h=None):
    if not os.path.exists(path): return None
    if h: return sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return sl.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def lbl(sl, x, y, text, dark=False):
    T(sl, x, y, 9, 0.24, text.upper(), sz=7, bold=True,
      c=LGREY if dark else STXT)

def hrule(sl, y, dark=False):
    R(sl, 0.55, y, W - 1.1, 0.012, MGREY if dark else RULEW)

def new(): return prs.slides.add_slide(BLANK)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE  (dark)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, BLACK)
lbl(sl, 0.6, 0.38, "00. Introduction", dark=True)
hrule(sl, 0.68, dark=True)

T(sl, 0.6, 1.0, 12.0, 2.6,
  "Predicting Professional\nCS2 Match Outcomes",
  sz=52, bold=True, c=WHITE)

T(sl, 0.6, 3.72, 12.0, 0.5,
  "A Comparative Study of Logistic Regression, Random Forest, and Gradient Boosting.",
  sz=14, italic=True, c=LGREY)

hrule(sl, 4.42, dark=True)

T(sl, 0.6, 4.62, 12.0, 0.42,
  "Seisenbek Dias   ·   Mergen Temirzhan   ·   Onashov Aidos",
  sz=16, bold=True, c=LTXT)
T(sl, 0.6, 5.12, 9.0, 0.35,
  "Mining signal from noise in 7,033 elite competitive matches.",
  sz=12, italic=True, c=LGREY)

cards = [
    ("Seisenbek Dias",   "Problem  ·  Data  ·  Variables"),
    ("Mergen Temirzhan", "Models  ·  Setup  ·  Results"),
    ("Onashov Aidos",    "Insights  ·  Discussion  ·  Conclusion"),
]
for i, (name, topic) in enumerate(cards):
    cx = 0.6 + i * 4.24
    R(sl, cx, 5.78, 4.0, 1.12, DGREY)
    R(sl, cx, 5.78, 0.04, 1.12, WHITE)
    T(sl, cx + 0.2, 5.92, 3.7, 0.42, name,  sz=13, bold=True, c=WHITE)
    T(sl, cx + 0.2, 6.42, 3.7, 0.36, topic, sz=10, c=LGREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM  (dark, bold statement)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, BLACK)
lbl(sl, 0.6, 0.38, "01. The Problem", dark=True)
hrule(sl, 0.68, dark=True)

T(sl, 0.6, 1.1, 12.2, 2.6,
  "Top-tier teams look\nidentical on paper.",
  sz=60, bold=True, c=WHITE)

hrule(sl, 4.0, dark=True)

T(sl, 0.6, 4.22, 12.0, 1.5,
  "At the highest levels of Counter-Strike 2, aggregate statistics lose their discriminatory power.\n"
  "Predicting a winner before the match starts is a problem of stochastic variation, not just historical data.",
  sz=15, c=LGREY)

T(sl, 0.6, 6.92, 5.0, 0.28, "Seisenbek Dias", sz=9, c=LGREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE DATA  (light, giant numbers)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, OFFWH)
lbl(sl, 0.6, 0.38, "02. The Evidence")
hrule(sl, 0.68)

for i, (num, desc) in enumerate([
    ("7,033", "Professional\nMatches"),
    ("648",   "Global\nTournaments"),
    ("18",    "Months of Data\n(May 2024 – Oct 2025)"),
]):
    cx = 0.6 + i * 4.24
    T(sl, cx, 0.88, 4.0, 2.15, num, sz=90, bold=True, c=DTXT)
    R(sl, cx, 2.88, 3.8, 0.045, DTXT)
    T(sl, cx, 3.08, 4.0, 0.95, desc, sz=16, c=STXT)

hrule(sl, 4.3)

T(sl, 0.6, 4.5, 7.5, 0.95,
  "Post-Transition Architecture: Most prior studies rely on CS:GO data. "
  "The 2023 engine transition and the HLTV 2.0 rating system effectively reset the field.\n"
  "This study processes only native CS2 data from the current competitive landscape.",
  sz=11, c=STXT)

T(sl, 0.6, 6.92, 5.0, 0.28, "Seisenbek Dias", sz=9, c=STXT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — THE VARIABLES  (light, left blocks + right heatmap)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, OFFWH)
lbl(sl, 0.6, 0.38, "03. The Variables")
hrule(sl, 0.68)

for i, (code, name, desc) in enumerate([
    ("X1", "Mean HLTV 2.0 Rating",    "The team's average aggregate performance score"),
    ("X2", "Mean Kills Per Round",     "The aggregate historical combat lethality"),
    ("X3", "Head-to-Head Win Rate",    "Specific historical matchup dominance between the two teams"),
    ("X4", "Map-Specific Win Rate",    "The team's historical success rate on the decided map"),
]):
    y = 0.88 + i * 1.38
    R(sl, 0.6, y, 0.58, 1.18, DTXT)
    T(sl, 0.6,  y + 0.3, 0.58, 0.56, code,
      sz=14, bold=True, c=WHITE, align=PP_ALIGN.CENTER)
    T(sl, 1.32, y + 0.08, 5.5, 0.46, name,  sz=16, bold=True, c=DTXT)
    T(sl, 1.32, y + 0.6,  5.5, 0.46, desc,  sz=11, c=STXT)
    if i < 3: R(sl, 0.6, y + 1.18, 6.35, 0.012, RULEW)

T(sl, 0.6, 6.6, 7.0, 0.36,
  "* Pearson correlations: all feature-target |r| < 0.15  —  weak individual signal",
  sz=9, italic=True, c=STXT)
T(sl, 0.6, 6.92, 4.0, 0.28, "Seisenbek Dias", sz=9, c=STXT)

I(sl, f"{FIG}/jfig6_corr.png", 7.25, 0.78, 5.7, 5.65)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — THE MODELS  (light, 3 columns + formula boxes)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, OFFWH)
lbl(sl, 0.6, 0.38, "04. The Models")
hrule(sl, 0.68)

for i, (name, sub, formula, desc) in enumerate([
    ("Logistic\nRegression",
     "Linear Baseline",
     "P(Y=1|X) = 1 / (1 + e^(-b0 - b1X1 - b2X2 - b3X3 - b4X4))",
     "Maximum likelihood with L2 regularisation.\nEvaluates standardised log-odds per standard deviation."),
    ("Random\nForest",
     "Bagging Ensemble",
     "P(Y=1|X) = (1 / 200)  x  SUM [Pb(Y=1|X)]",
     "200 decision trees aggregating predictions by majority vote.  Robust to overfitting."),
    ("Gradient\nBoosting",
     "Additive Ensemble",
     "Fm(x) = Fm-1(x)  +  n * hm(x)    [n = 0.05,  M = 200]",
     "Sequential shallow trees fitting pseudo-residuals.\nPrioritises hold-out accuracy."),
]):
    cx = 0.6 + i * 4.24
    R(sl, cx, 0.88, 4.0, 5.92, CARD)
    R(sl, cx, 0.88, 4.0, 0.055, DTXT)
    T(sl, cx + 0.2, 1.04, 3.6, 0.95, name,    sz=22, bold=True, c=DTXT)
    T(sl, cx + 0.2, 2.05, 3.6, 0.34, sub,     sz=11, italic=True, c=STXT)
    R(sl, cx + 0.2, 2.48, 3.6, 0.012, RULEW)
    # Formula box
    R(sl, cx + 0.2, 2.62, 3.6, 1.05, DTXT)
    T(sl, cx + 0.22, 2.72, 3.56, 0.85, formula,
      sz=9, bold=True, c=WHITE, align=PP_ALIGN.CENTER)
    R(sl, cx + 0.2, 3.77, 3.6, 0.012, RULEW)
    T(sl, cx + 0.2, 3.92, 3.6, 1.6, desc, sz=11, italic=True, c=STXT)

T(sl, 0.6, 6.92, 5.0, 0.28, "Mergen Temirzhan", sz=9, c=STXT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — THE VERDICT  (light, 3 result blocks)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, OFFWH)
lbl(sl, 0.6, 0.38, "05. The Verdict")
hrule(sl, 0.68)

for i, (name, acc, f1, auc, best) in enumerate([
    ("Logistic\nRegression", "56.1%", "F1: 0.484", "AUC: 0.574", False),
    ("Random\nForest",       "57.2%", "F1: 0.536", "AUC: 0.571", True),
    ("Gradient\nBoosting",   "56.4%", "F1: 0.528", "AUC: 0.573", False),
]):
    cx = 0.6 + i * 4.24
    bg = DTXT if best else CARD
    pt = WHITE if best else DTXT
    st = LGREY if best else STXT
    rl = MGREY if best else RULEW
    R(sl, cx, 0.88, 4.0, 5.55, bg)
    T(sl, cx + 0.22, 1.02, 3.56, 0.95, name, sz=20, bold=True, c=pt)
    R(sl, cx + 0.22, 2.08, 3.56, 0.018, rl)
    T(sl, cx + 0.22, 2.22, 3.56, 1.45, acc,
      sz=56, bold=True, c=pt, align=PP_ALIGN.CENTER)
    T(sl, cx + 0.22, 3.78, 3.56, 0.36, "Accuracy",
      sz=11, c=st, align=PP_ALIGN.CENTER)
    R(sl, cx + 0.22, 4.28, 3.56, 0.018, rl)
    T(sl, cx + 0.22, 4.5,  3.56, 0.4,  f1,
      sz=17, bold=True, c=pt, align=PP_ALIGN.CENTER)
    T(sl, cx + 0.22, 4.95, 3.56, 0.4,  auc,
      sz=17, c=st, align=PP_ALIGN.CENTER)

T(sl, 0.6, 6.56, W - 1.2, 0.36,
  "Paired t-tests (5-fold CV) confirm: differences between classifiers are not statistically significant (p > 0.05)",
  sz=10, italic=True, c=STXT, align=PP_ALIGN.CENTER)
T(sl, 0.6, 6.92, 5.0, 0.28, "Mergen Temirzhan", sz=9, c=STXT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — 57.2% vs 54.6%  (dark, giant numbers)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, BLACK)
lbl(sl, 0.6, 0.38, "06. The Baseline", dark=True)
hrule(sl, 0.68, dark=True)

T(sl, 0.6, 0.95, 5.6, 1.1,
  "The Bound of\nStochastic Variation",
  sz=24, bold=True, c=WHITE)
T(sl, 0.6, 2.22, 5.6, 2.8,
  "A naïve baseline — simply predicting the majority class "
  "(always predict 'Win') — yields an automatic 54.6% accuracy.\n\n"
  "Applying a 200-tree Random Forest model lifts that accuracy "
  "by just 2.6%.\n\nThe models are not failing; they have reached "
  "the limits of pre-match data.",
  sz=12, c=LGREY)

T(sl, 6.9, 0.72, 6.0, 2.1, "57.2%", sz=92, bold=True, c=WHITE)
R(sl, 6.9, 2.95, 5.9, 0.042, MGREY)
T(sl, 6.9, 3.12, 6.0, 0.5,
  "Random Forest Maximum Accuracy", sz=14, c=LGREY)

T(sl, 6.9, 4.05, 6.0, 2.0, "54.6%", sz=92, bold=True, c=LGREY)
R(sl, 6.9, 6.08, 5.9, 0.042, MGREY)
T(sl, 6.9, 6.28, 6.0, 0.42, "Naïve Baseline", sz=14, c=LGREY)

T(sl, 0.6, 6.92, 5.0, 0.28, "Mergen Temirzhan", sz=9, c=LGREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — CONFUSION MATRICES  (light, full figure)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, OFFWH)
lbl(sl, 0.6, 0.38, "07. Confusion Matrices")
hrule(sl, 0.68)

I(sl, f"{FIG}/jfig1_cm.png", 0.6, 0.85, 12.1, 5.75)

T(sl, 0.6, 6.68, W - 1.2, 0.34,
  "All models over-predict Win — false positive rate is the dominant error, as expected from weak feature-target correlations",
  sz=10, italic=True, c=STXT, align=PP_ALIGN.CENTER)
T(sl, 0.6, 6.92, 5.0, 0.28, "Mergen Temirzhan", sz=9, c=STXT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — THE SIGNAL  (dark, 01/02/03 large numbers)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, BLACK)
lbl(sl, 0.6, 0.38, "08. The Signal", dark=True)
hrule(sl, 0.68, dark=True)

for i, (num, name, imp, desc) in enumerate([
    ("01", "Map Win Rate",
     "Importance: 0.403 (RF)  ·  0.399 (GB)",
     "The dominant predictor. At the elite tier, being pulled into an unfamiliar "
     "structural environment — a weak map — creates a measurable, objective disadvantage "
     "that raw skill cannot overcome."),
    ("02", "Head-to-Head",
     "Importance: 0.269 (RF)  ·  0.292 (GB)",
     "The secondary signal. Stylistic matchups and historical dominance persist over "
     "time, creating psychological and tactical friction between opponents."),
    ("03", "Aggregate Skill",
     "Importance: 0.191 HLTV  ·  0.138 KPR",
     "The weakest signal. When both teams are elite, historical combat lethality "
     "and overall ratings cancel each other out."),
]):
    cx = 0.6 + i * 4.24
    T(sl, cx, 0.88, 3.9, 1.6, num, sz=76, bold=True, c=DNUM)
    T(sl, cx, 2.3,  3.9, 0.55, name, sz=21, bold=True, c=WHITE)
    T(sl, cx, 2.92, 3.9, 0.36, imp,  sz=11, bold=True, c=LGREY)
    R(sl, cx, 3.38, 3.9, 0.012, MGREY)
    T(sl, cx, 3.52, 3.9, 2.65, desc, sz=11, c=LGREY)

T(sl, 0.6, 6.92, 5.0, 0.28, "Onashov Aidos", sz=9, c=LGREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — FEATURE IMPORTANCE FIGURE  (light, full chart)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, OFFWH)
lbl(sl, 0.6, 0.38, "09. Feature Importance — Both Ensemble Models")
hrule(sl, 0.68)

I(sl, f"{FIG}/jfig3_imp.png", 0.6, 0.85, 12.1, 5.75)

T(sl, 0.6, 6.68, W - 1.2, 0.34,
  "Map Win Rate ranks first in both RF and GB with ~40% importance. KPR is consistently the weakest contributor.",
  sz=10, italic=True, c=STXT, align=PP_ALIGN.CENTER)
T(sl, 0.6, 6.92, 5.0, 0.28, "Onashov Aidos", sz=9, c=STXT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — THE FEATURE CEILING  (dark, pre-match vs in-game)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, BLACK)
lbl(sl, 0.6, 0.38, "10. The Feature Ceiling", dark=True)
hrule(sl, 0.68, dark=True)

R(sl, W/2 - 0.022, 0.85, 0.044, 5.75, MGREY)

T(sl, 0.6, 0.92, 6.0, 0.6, "Pre-Match Horizon",
  sz=22, bold=True, c=WHITE)
for j, item in enumerate(["Historical Aggregates", "Temporal Noise",
                           "Map Veto", "Roster Shifts"]):
    y = 1.7 + j * 0.78
    R(sl, 0.6, y, 5.8, 0.6, DGREY)
    T(sl, 0.82, y + 0.12, 5.4, 0.36, item, sz=15, c=LTXT)
R(sl, 0.6, 5.6, 5.8, 0.72, RGBColor(0x16, 0x16, 0x16))
T(sl, 0.8, 5.68, 5.4, 0.5, "Max Accuracy:  ~56%",
  sz=19, bold=True, c=WHITE)

T(sl, 7.12, 0.92, 5.9, 0.6, "In-Game Reality",
  sz=22, bold=True, c=WHITE)
for j, item in enumerate(["Real-Time Economy", "Alive-Player State",
                           "Round Score", "Clutch Variance"]):
    y = 1.7 + j * 0.78
    R(sl, 7.12, y, 5.8, 0.6, DGREY)
    T(sl, 7.34, y + 0.12, 5.4, 0.36, item, sz=15, c=LTXT)
R(sl, 7.12, 5.6, 5.8, 0.72, RGBColor(0x16, 0x16, 0x16))
T(sl, 7.32, 5.68, 5.4, 0.5, "Max Accuracy:  80%+",
  sz=19, bold=True, c=WHITE)

T(sl, 0.6, 6.92, 5.0, 0.28, "Onashov Aidos", sz=9, c=LGREY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION  (dark, single memorable quote + findings)
# ═══════════════════════════════════════════════════════════════════
sl = new()
R(sl, 0, 0, W, H, BLACK)
lbl(sl, 0.6, 0.38, "11. Conclusion", dark=True)
hrule(sl, 0.68, dark=True)

T(sl, 0.6, 1.05, 12.2, 0.78,
  "The algorithms do not lack capacity.",
  sz=38, bold=True, c=WHITE)
T(sl, 0.6, 1.9, 12.2, 0.78,
  "The data lacks certainty.",
  sz=38, bold=True, c=LGREY)

hrule(sl, 2.95, dark=True)

T(sl, 0.6, 3.15, 12.2, 1.2,
  "At the highest tier of competition, pre-match outcomes are bound by a structural ceiling of stochastic variance.\n"
  "To break the 56% barrier, we must look past aggregate history and measure the game in motion.",
  sz=14, c=LGREY)

for i, (h, verdict, finding) in enumerate([
    ("H3", "confirmed",       "Map Win Rate is the dominant predictor"),
    ("H2", "partial",         "H2H Win Rate is the secondary signal"),
    ("H1", "weakly supported","HLTV Rating / KPR — least discriminating"),
]):
    y = 4.52 + i * 0.42
    T(sl, 0.6,  y, 0.5, 0.36, h,       sz=11, bold=True, c=LGREY)
    T(sl, 1.1,  y, 2.2, 0.36, verdict, sz=11, italic=True, c=LGREY)
    T(sl, 3.3,  y, 9.0, 0.36, finding, sz=11, c=RGBColor(0x55, 0x55, 0x55))

hrule(sl, 5.88, dark=True)
T(sl, 0.6, 6.06, W - 1.2, 0.38,
  "github.com/aloneen/cs2-match-outcome-prediction",
  sz=11, c=LGREY, align=PP_ALIGN.CENTER)
T(sl, 0.6, 6.55, W - 1.2, 0.6,
  "Thank you  ·  Questions?",
  sz=30, bold=True, c=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
OUT = "final/defense_presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({os.path.getsize(OUT)//1024} KB)  |  12 slides")
